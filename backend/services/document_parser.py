import os
import pandas as pd
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from typing import List

def extract_text_from_pdf(file_path: str) -> str:
    """Extracts text content from a PDF file using PyMuPDF."""
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text() + "\n"
    except Exception as e:
        raise ValueError(f"Error reading PDF file: {str(e)}")
    return text

def extract_text_from_docx(file_path: str) -> str:
    """Extracts paragraph and table text content from a DOCX file."""
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                text += " | ".join(row_text) + "\n"
    except Exception as e:
        raise ValueError(f"Error reading DOCX file: {str(e)}")
    return text

def extract_text_from_pptx(file_path: str) -> str:
    """Extracts text content from slide shapes in a PPTX file."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
    except Exception as e:
        raise ValueError(f"Error reading PPTX file: {str(e)}")
    return text

def extract_text_from_excel(file_path: str) -> str:
    """Extracts sheets from an Excel file, converting them to markdown tables."""
    text = ""
    try:
        xl = pd.ExcelFile(file_path)
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            text += f"Sheet: {sheet_name}\n"
            text += df.to_markdown(index=False) + "\n\n"
    except Exception as e:
        try:
            df = pd.read_excel(file_path)
            text += df.to_string(index=False)
        except Exception as e2:
            raise ValueError(f"Error reading Excel file: {str(e2)}")
    return text

def extract_text_from_csv(file_path: str) -> str:
    """Converts CSV table records to markdown table representation."""
    try:
        df = pd.read_csv(file_path)
        return df.to_markdown(index=False)
    except Exception as e:
        raise ValueError(f"Error reading CSV file: {str(e)}")

def extract_text_from_txt(file_path: str) -> str:
    """Extracts text from a plain TXT or markdown file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        raise ValueError(f"Error reading TXT file: {str(e)}")

def extract_text(file_path: str, filename: str) -> str:
    """
    Main entry point for document text extraction.
    Ensures file-specific parser wrapping and error containment.
    """
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in [".docx", ".doc"]:
        return extract_text_from_docx(file_path)
    elif ext in [".pptx", ".ppt"]:
        return extract_text_from_pptx(file_path)
    elif ext in [".xlsx", ".xls"]:
        return extract_text_from_excel(file_path)
    elif ext == ".csv":
        return extract_text_from_csv(file_path)
    elif ext in [".txt", ".md"]:
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")

def chunk_text(text: str, chunk_size: int = 400, overlap: int = 50) -> List[str]:
    """
    Splits text into semantic chunks of roughly chunk_size words, 
    with a designated context overlap. Respects paragraph spacing where possible.
    """
    paragraphs = text.split("\n")
    chunks = []
    current_chunk = []
    current_word_count = 0

    for para in paragraphs:
        if not para.strip():
            continue
        words = para.split()
        if current_word_count + len(words) > chunk_size and current_chunk:
            chunks.append(" ".join(current_chunk))
            # Capture overlap text
            overlap_words = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
            current_chunk = list(overlap_words)
            current_word_count = sum(len(w.split()) for w in current_chunk)
        
        current_chunk.append(para)
        current_word_count += len(words)

    if current_chunk:
        chunks.append(" ".join(current_chunk))
    
    return chunks
